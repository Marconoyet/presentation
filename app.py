import sys  # nopep8
sys.path.insert(0, './lib')  # nopep8

import requests
import json
from flask import Flask, request, send_file
from flask_cors import CORS
from pptx import Presentation
import os
import io
import uuid

from helpers.ppt_helpers import replace_text_in_shapes, clone_post_slide
from helpers.data_utils import extract_text_and_image
app = Flask(__name__)
CORS(app)
TEMPLATE_PATH = "templates/base.pptx"
OUTPUT_DIR = "generated"


@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    data = request.json or {}
    project_id = data.get("project_id")
    tasklist_id = data.get("tasklist_id")
    generate_type = data.get("generate_type")

    # إعداد بيانات Zoho
    zoho_url = "https://www.zohoapis.com/crm/v7/functions/projectpresentation/actions/execute"
    headers = {"Content-Type": "application/json"}
    payload = {
        "type": generate_type,
        "data": {
            "projectId": project_id,
            "tasklistId": tasklist_id
        }
    }
    params = {
        "auth_type": "apikey",
        "zapikey": "1003.3c9532ce6454d6df1b898e83bfff40ea.bed6415c9b3cf64cc4532c1755fee909",
    }

    try:
        response = requests.post(
            zoho_url, headers=headers, json=payload, params=params)
        response.raise_for_status()
        zoho_result = response.json()

        # ✅ استخراج output وتحويله إلى dict
        raw_output = zoho_result.get("details", {}).get("output", "{}")
        parsed_output = json.loads(raw_output)

        posts = parsed_output.get("posts", [])
        project_name = parsed_output.get("projectName", "Unnamed Project")
        calendar_name = parsed_output.get(
            "taskCalenderName", "No Calendar Name")

    except requests.exceptions.RequestException as e:
        return {"error": f"Zoho API call failed: {str(e)}"}, 500
    except json.JSONDecodeError as e:
        return {"error": f"Failed to parse Zoho output: {str(e)}"}, 500

    # ✅ إنشاء العرض التقديمي
    prs = Presentation(TEMPLATE_PATH)

    # الشريحة الأولى
    replace_text_in_shapes(prs.slides[0].shapes, {
        "project_name": project_name,
        "calendar_name": calendar_name
    })

    post_template_slide = prs.slides[1]
    end_slide_template = prs.slides[2]

    # ✅ توليد شرائح من البيانات القادمة من Zoho
    for index, post in enumerate(posts, start=1):
        description = post.get("description", "")
        caption, image = extract_text_and_image(description)

        clone_post_slide(prs, post_template_slide, {
            "caption": caption,
            "n": str(index),
            "image": image
        })

    clone_post_slide(prs, end_slide_template, {})

    # حذف الشرائح الأصلية
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[2])
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])

    # حفظ وتصدير الملف
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)  # إعادة المؤشر للبداية

    # ✅ إرسال الملف مباشرة دون حفظه
    return send_file(
        pptx_stream,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=f"{project_name}.pptx"
    )


@app.route('/sdk')
def sdk():
    try:
        import zcatalyst_sdk as zcatalyst
        app = zcatalyst.initialize(req=request)
        cache_resp = app.cache().segment().put('Key', 'value')
        return cache_resp, 200
    except Exception as e:
        return 'Got exception: ' + repr(e)


if __name__ == '__main__':
    listen_port = int(os.getenv('X_ZOHO_CATALYST_LISTEN_PORT', 9000))
    app.run(host="0.0.0.0", port=listen_port)
