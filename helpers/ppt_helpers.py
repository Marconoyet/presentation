import sys  # nopep8
sys.path.insert(0, './lib')  # nopep8

import re
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.util import Pt
from copy import deepcopy
import os
import requests
from bs4 import BeautifulSoup
from io import BytesIO
from PIL import Image

from pptx.util import Pt

MAX_LINES = 15
MAX_CHARS_PER_LINE = 36
DEFAULT_FONT_SIZE = Pt(18)
MIN_FONT_SIZE = Pt(8)


def fit_text_by_box_ratio(shape):
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text = shape.text

    lines = calculate_required_lines(text)
    if lines <= MAX_LINES:
        final_font = DEFAULT_FONT_SIZE
    else:
        ratio = MAX_LINES / lines
        scaled_size = DEFAULT_FONT_SIZE.pt * ratio
        final_font = Pt(max(scaled_size, MIN_FONT_SIZE.pt))

    # Apply font
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = final_font


def calculate_required_lines(text):
    # يقسم بناءً على فواصل الأسطر الفعلية
    raw_lines = text.splitlines()

    total_lines = 0
    for line in raw_lines:
        # نحسب كم سطر فعلي محتاجه كل سطر بناءً على عدد الحروف
        line_length = len(line)
        if line_length == 0:
            total_lines += 1  # سطر فاضي
        else:
            # لو السطر طوله أكبر من الحد، نقسمه لعدة أسطر
            lines_needed = (line_length - 1) // MAX_CHARS_PER_LINE + 1
            total_lines += lines_needed

    return total_lines


def is_arabic_line(text):
    arabic_chars = re.findall(r'[\u0600-\u06FF]', text)
    return len(arabic_chars)


def is_latin_line(text):
    latin_chars = re.findall(r'[A-Za-z]', text)
    return len(latin_chars)


def replace_text_in_shapes(shapes, replacements, font_color=None):
    for shape in shapes:
        if shape.has_text_frame:
            full_text = shape.text
            is_caption = False
            caption_value = ""

            for key, value in replacements.items():
                if key == "image":
                    continue
                placeholder = f"{{{{{key}}}}}"
                if placeholder in full_text:
                    full_text = full_text.replace(placeholder, str(value))
                    if key == "caption":
                        is_caption = True
                        caption_value = str(value)

            if is_caption:
                text_frame = shape.text_frame
                text_frame.clear()  # احذف الفقرات القديمة

                lines = caption_value.splitlines()
                for i, line in enumerate(lines):
                    has_colon = ':' in line
                    contains_arabic = is_arabic_line(line.strip()) > 0
                    if has_colon and contains_arabic:
                        parts = line.split(':', 1)
                        line = f"{parts[1].strip()} :{parts[0].strip()}"
                    paragraph = text_frame.add_paragraph(
                    ) if i > 0 else text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    link_match = re.search(r'(https?://\S+)', line)
                    if link_match:
                        url = link_match.group(1)
                        clean_text = line.replace(url, '').strip()
                        run.text = clean_text + ' '
                        hlink_run = paragraph.add_run()
                        hlink_run.text = url
                        hlink_run.hyperlink.address = url

                    else:
                        run.text = line

                    if has_colon and contains_arabic:
                        # ✅ لو فيه ":" → اتجاه من اليسار لليمين لكن المحاذاة يمين
                        paragraph._pPr.set('algn', 'r')  # Align right
                        paragraph._pPr.set('rtl', '0')  # LTR direction
                    else:
                        # الافتراضي حسب اللغة
                        if is_arabic_line(line) > is_latin_line(line):
                            paragraph._pPr.set('algn', 'r')
                            paragraph._pPr.set('rtl', '1')
                        else:
                            paragraph._pPr.set('algn', 'l')
                            paragraph._pPr.set('rtl', '0')

                    if font_color:
                        run.font.color.rgb = font_color

                fit_text_by_box_ratio(shape)
            else:
                # 👇 باقي العناصر
                shape.text = full_text
            if font_color:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = font_color


def clone_post_slide(prs, template_slide, replacements=None):
    if replacements is None:
        replacements = {}

    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # انسخ العناصر
    for shape in template_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # استبدال النصوص بلون أبيض
    white = RGBColor(255, 255, 255)
    replace_text_in_shapes(new_slide.shapes, replacements, font_color=white)

    # لو فيه صورة
    if "image" in replacements and replacements["image"]:
        image_url = replacements["image"]

        # ابحث عن العنصر اللي فيه {{image}}
        image_placeholder = None
        for shape in new_slide.shapes:
            if shape.has_text_frame and "{{image}}" in shape.text:
                image_placeholder = shape
                break
        if image_placeholder:
            # احصل على حجم ومكان الإطار (الـ TextBox القديم)
            box_left = image_placeholder.left
            box_top = image_placeholder.top
            box_width = image_placeholder.width
            box_height = image_placeholder.height

            # احذف صندوق النص
            new_slide.shapes._spTree.remove(image_placeholder._element)
            try:
                # حمل الصورة من الإنترنت
                response = requests.get(image_url, timeout=10)
                response.raise_for_status()
                tmp_path = "temp_image.png"
                with Image.open(BytesIO(response.content)) as image:
                    image.save(tmp_path)

                    # أبعاد الصورة الأصلية بالبكسل
                    img_width_px, img_height_px = image.size

                    # نحول حجم البوكس من EMU إلى بكسل تقريبي (بافتراض 96 DPI)
                    EMU_PER_INCH = 914400
                    DPI = 96
                    def emu_to_px(emu): return int(emu * DPI / EMU_PER_INCH)

                    box_width_px = emu_to_px(box_width)
                    box_height_px = emu_to_px(box_height)

                    # احسب النسبة بين البوكس والصورة
                    width_ratio = box_width_px / img_width_px
                    height_ratio = box_height_px / img_height_px
                    scale_ratio = min(width_ratio, height_ratio)

                    # الأبعاد الجديدة بعد التصغير
                    final_width_px = int(img_width_px * scale_ratio)
                    final_height_px = int(img_height_px * scale_ratio)

                    def px_to_emu(px): return int(px * EMU_PER_INCH / DPI)
                    final_width = px_to_emu(final_width_px)
                    final_height = px_to_emu(final_height_px)

                    final_left = box_left + int((box_width - final_width) / 2)
                    final_top = box_top + int((box_height - final_height) / 2)

                # ✅ بعد ما الصورة تقفلت تمامًا
                new_slide.shapes.add_picture(
                    tmp_path, final_left, final_top,
                    width=final_width, height=final_height
                )

            except Exception as e:
                print(f"Error adding image: {e}")
            return new_slide
